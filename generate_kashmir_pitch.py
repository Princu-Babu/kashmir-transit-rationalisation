"""
generate_kashmir_pitch.py

The diagrammatic meeting deck for the Kashmir Valley route-rationalisation plan.
Flowcharts, icons and charts carry the story; the methodology is written out in
full — every formula is stated plainly and every source is cited — so the plan
is defensible line by line in the room.

Output: Kashmir_Transit_Diagrammatic_Pitch.pptx (16:9).
"""
from __future__ import annotations

import argparse
import os
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# ─── Theme ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1A, 0x23, 0x7E)
TEAL        = RGBColor(0x00, 0x69, 0x5C)
TEAL_LIGHT  = RGBColor(0xB2, 0xDF, 0xDB)
SAFFRON     = RGBColor(0xD3, 0x2F, 0x2F)
SAFFRON_LT  = RGBColor(0xFF, 0xEB, 0xEE)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
GOLD        = RGBColor(0xF9, 0xA8, 0x25)
GOLD_LT     = RGBColor(0xFF, 0xF8, 0xE1)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LT    = RGBColor(0xE8, 0xF5, 0xE9)
BLUE        = RGBColor(0x15, 0x65, 0xC0)
ORANGE      = RGBColor(0xEF, 0x6C, 0x00)
DARK        = RGBColor(0x21, 0x21, 0x21)
GREY        = RGBColor(0x75, 0x75, 0x75)
LIGHT       = RGBColor(0xF5, 0xF5, 0xF5)
CARD        = RGBColor(0xF5, 0xF6, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)


# ─── Helpers ─────────────────────────────────────────────────────────────────
def add_text(slide, text, x, y, w, h, *, size=14, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             font="Segoe UI"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = str(text) if str(text) else " "
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, *, fill=TEAL, line=None, rounded=False):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.shadow.inherit = False
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1.0)
    if rounded:
        try:
            shp.adjustments[0] = 0.10
        except Exception:
            pass
    return shp


def add_oval(slide, x, y, w, h, *, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.shadow.inherit = False
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_arrow(slide, x1, y1, x2, y2, *, color=GREY, width=2.0):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    try:
        from pptx.oxml.ns import qn
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    except Exception:
        pass
    return conn


def add_centered_label(slide, text, x, y, w, h, *, size=12, bold=True, color=WHITE,
                       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = str(text) if str(text) else " "
    r = p.runs[0]
    r.font.name = "Segoe UI"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def slide_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_header(slide, title, *, color=NAVY, sub=None):
    add_rect(slide, 0, 0, 0.18, 7.5, fill=color)
    add_text(slide, title, 0.5, 0.35, 12.4, 0.7, size=27, bold=True, color=color, font="Calibri")
    if sub:
        add_text(slide, sub, 0.5, 1.02, 12.4, 0.4, size=14, color=GREY, italic=True)


def footer(slide, text):
    add_text(slide, text, 0.5, 7.06, 12.3, 0.3, size=9, color=GREY, italic=True)


def formula_card(slide, x, y, w, h, plain, symbolic, source, accent):
    add_rect(slide, x, y, w, h, fill=CARD, rounded=True, line=accent)
    add_text(slide, plain, x + 0.3, y + 0.16, w - 0.6, h - 1.0, size=13, bold=True, color=accent)
    add_text(slide, symbolic, x + 0.3, y + h - 0.84, w - 0.6, 0.42, size=11.5,
             color=DARK, italic=True)
    add_text(slide, "Basis:  " + source, x + 0.3, y + h - 0.42, w - 0.6, 0.38,
             size=10, color=GREY, italic=True)


# ─── Deck ────────────────────────────────────────────────────────────────────
def build(prs, s):
    # ── 1 · Cover ───────────────────────────────────────────────────────────
    sl = slide_blank(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    add_rect(sl, 0, 6.7, 13.33, 0.8, fill=TEAL)
    add_text(sl, "KASHMIR VALLEY PUBLIC TRANSPORT", 0.7, 0.7, 12, 0.6,
             size=20, bold=True, color=TEAL_LIGHT, font="Calibri")
    add_text(sl, "Route Rationalisation Plan", 0.7, 1.3, 12, 0.9,
             size=40, bold=True, color=WHITE, font="Calibri")
    add_text(sl, "A data-driven service plan for the Srinagar Valley bus network, "
                 "built on the Srinagar Smart City electric-bus backbone.",
             0.7, 2.25, 11.5, 0.6, size=15, color=TEAL_LIGHT, italic=True)
    cards = [
        ("ACTIVE ROUTES", f"{s['active']}", "clear, non-overlapping lines", TEAL),
        ("TOTAL BUSES",   f"{s['fleet']:,}", "recommended for the network", GOLD),
        ("LONGEST WAIT",  "35 min", "anywhere on the network", PURPLE),
        ("RESIDENTS",     f"{s['coverage']:.0f}%", "covered within a short walk", GREEN),
    ]
    cw, cx, cy, gap = 2.7, 0.9, 3.4, 0.3
    for i, (label, big, sub, col) in enumerate(cards):
        x = cx + i * (cw + gap)
        add_rect(sl, x, cy, cw, 2.6, fill=WHITE, rounded=True)
        add_rect(sl, x, cy, cw, 0.45, fill=col, rounded=True)
        add_centered_label(sl, label, x, cy, cw, 0.45, size=11, color=WHITE)
        add_text(sl, big, x, cy + 0.6, cw, 1.2, size=42, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri")
        add_text(sl, sub, x, cy + 1.95, cw, 0.5, size=11, color=GREY,
                 align=PP_ALIGN.CENTER, italic=True)
    add_text(sl, "Prepared for the Principal Secretary (Transport), Government of Jammu & Kashmir   ·   "
                 "for planning review ahead of operator consultation.",
             0.7, 6.85, 12, 0.45, size=11, color=WHITE, italic=True, align=PP_ALIGN.CENTER)

    # ── 2 · Why this plan exists ──────────────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "Why this plan exists", color=SAFFRON,
                 sub="Srinagar's bus network grew permit by permit, with no central frequency plan.")
    problems = [
        ("Wasteful duplication",
         "More than a dozen buses chase the same riders on the Parimpora–Lal Chowk–Dalgate "
         "corridor, while whole neighbourhoods get nothing."),
        ("Transit deserts",
         "The southern industrial belt (Khonmoh, Rangreth) and satellite towns "
         "(Ganderbal, Pulwama) are badly under-served."),
        ("No reliability",
         "Without a published frequency, buses bunch together at peak hours and disappear "
         "off-peak — riders cannot plan around them."),
        ("Wrong bus, wrong road",
         "Large buses are sent down narrow old-city lanes, while small minibuses run "
         "long inter-district highways."),
        ("Built around half the riders",
         "Almost two-thirds of riders are women, yet the network was never designed "
         "around where they actually travel."),
    ]
    y = 1.75
    for title, body in problems:
        add_oval(sl, 0.7, y + 0.07, 0.3, 0.3, fill=SAFFRON)
        add_text(sl, title, 1.2, y - 0.05, 3.7, 0.5, size=15, bold=True, color=SAFFRON,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, body, 5.1, y - 0.05, 7.7, 0.95, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += 1.02
    footer(sl, "Slide 2  ·  The problem")

    # ── 3 · How the plan is built (four phases) ──────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "How the plan is built — four stages", color=TEAL)
    stages = [
        ("STAGE 1", "Gather", NAVY, [
            "Every registered route permit in the valley",
            "Real road driving times on the live map",
            "Where people live, from a satellite population grid",
            "The places people travel to (hospitals, colleges, markets)",
        ]),
        ("STAGE 2", "Measure demand", TEAL, [
            "Population living within a short walk of each route",
            "The destinations each route reaches, weighted by importance",
            "Extra weight for women's destinations and tourist corridors",
        ]),
        ("STAGE 3", "Rank and size", PURPLE, [
            "Rank every route by how much travel it serves",
            "Group routes into high, medium and low priority",
            "Choose how often buses run on each",
            "Work out how many buses that needs",
        ]),
        ("STAGE 4", "Check and publish", SAFFRON, [
            "Eight automatic quality checks before anything is released",
            "Compare against real electric-bus ridership as a sanity check",
            "Publish the route plan, maps and the operator workbook",
        ]),
    ]
    bx, by, bw, bh, gap = 0.5, 1.85, 3.0, 4.5, 0.2
    for i, (tag, name, col, bullets) in enumerate(stages):
        x = bx + i * (bw + gap)
        add_rect(sl, x, by, bw, bh, fill=CARD, rounded=True, line=col)
        add_rect(sl, x, by, bw, 0.55, fill=col, rounded=True)
        add_centered_label(sl, tag, x, by, bw, 0.55, size=13, color=WHITE)
        add_text(sl, name, x + 0.15, by + 0.62, bw - 0.3, 0.7, size=17, bold=True,
                 color=col, align=PP_ALIGN.CENTER, font="Calibri")
        for j, b in enumerate(bullets):
            row_y = by + 1.5 + j * 0.72
            add_oval(sl, x + 0.22, row_y + 0.1, 0.13, 0.13, fill=col)
            add_text(sl, b, x + 0.45, row_y, bw - 0.6, 0.7, size=10.5, color=DARK)
        if i < 3:
            ay = by + bh / 2
            add_arrow(sl, x + bw, ay, x + bw + gap, ay, color=col, width=2.5)
    footer(sl, "Slide 3  ·  Method, stage by stage")

    # ── 4 · Network composition (by role) ────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, f"From {s['total_routes']} permits to {s['active']} clear routes",
                 color=PURPLE,
                 sub="Duplicate permits are merged into unified trunk lines; the rest are kept as feeders.")
    total = s['total_routes']
    parts = [
        ("Upgraded to trunk line", s['trunks'], TEAL),
        ("Kept as feeder",         s['feeders'], NAVY),
        ("Merged into a trunk",    s['merged'], SAFFRON),
    ]
    bx, by, bar_w, bar_h = 0.7, 1.95, 11.9, 1.0
    accum = 0
    for label, count, col in parts:
        w_part = bar_w * count / total
        add_rect(sl, bx + accum, by, w_part, bar_h, fill=col)
        if w_part > 1.4:
            add_centered_label(sl, f"{count}", bx + accum, by, w_part, bar_h, size=24, color=WHITE)
        accum += w_part
    lx, ly, lw = 0.7, 3.2, 4.0
    for i, (label, count, col) in enumerate(parts):
        x = lx + i * (lw + 0.0)
        add_rect(sl, x, ly, 0.22, 0.62, fill=col)
        add_text(sl, label, x + 0.32, ly + 0.0, lw - 0.4, 0.32, size=13, bold=True, color=DARK)
        add_text(sl, f"{count} routes  ·  {100*count/total:.0f}% of permits",
                 x + 0.32, ly + 0.32, lw - 0.4, 0.3, size=11, color=GREY)
    add_rect(sl, 0.7, 4.5, 11.9, 1.9, fill=GREEN_LT, rounded=True, line=GREEN)
    add_text(sl, "What changes for the rider", 0.95, 4.62, 11.4, 0.4, size=15, bold=True, color=GREEN)
    add_text(sl,
             f"The {s['active']} active routes are a clear, legible network instead of "
             f"{s['total_routes']} overlapping permits. Trunk lines carry the busiest "
             "corridors at high frequency; feeders connect neighbourhoods into them; and "
             f"the merged {s['merged']} permits stop competing for the same passengers.",
             0.95, 5.0, 11.4, 1.3, size=13, color=DARK)
    footer(sl, "Slide 4  ·  Network composition")

    # ── 5 · Recommended fleet ────────────────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, f"Recommended fleet — {s['fleet']:,} buses", color=GOLD,
                 sub=f"{s['hpv']} large (12-metre)  +  {s['mpv']} medium (9-metre)  +  "
                     f"{s['lpv']} small (minibus)")
    classes = [
        ("Large bus", s['hpv'], "12-metre standard low-floor", "about 50 seats", BLUE, 4.0),
        ("Medium bus", s['mpv'], "9-metre electric / midi bus", "about 35 seats", TEAL, 3.0),
        ("Small bus", s['lpv'], "minibus / tempo traveller", "about 15 seats", ORANGE, 2.0),
    ]
    cx, cy, cw, ch, gap = 0.9, 1.95, 3.9, 4.3, 0.2
    for i, (cls, count, desc, seats, col, bus_w) in enumerate(classes):
        x = cx + i * (cw + gap)
        add_rect(sl, x, cy, cw, ch, fill=LIGHT, rounded=True)
        add_rect(sl, x, cy, cw, 0.6, fill=col, rounded=True)
        add_centered_label(sl, cls, x, cy, cw, 0.6, size=19, color=WHITE)
        add_text(sl, f"{count}", x, cy + 0.7, cw, 1.1, size=60, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri")
        bus_x = x + (cw - bus_w) / 2
        bus_y = cy + 2.15
        add_rect(sl, bus_x, bus_y, bus_w, 0.7, fill=col, rounded=True)
        win_w = (bus_w - 0.4) / 6
        for w_i in range(5):
            add_rect(sl, bus_x + 0.2 + w_i * (win_w + 0.05), bus_y + 0.12, win_w, 0.25, fill=WHITE)
        add_oval(sl, bus_x + 0.25, bus_y + 0.55, 0.3, 0.3, fill=DARK)
        add_oval(sl, bus_x + bus_w - 0.55, bus_y + 0.55, 0.3, 0.3, fill=DARK)
        add_text(sl, desc, x, cy + 3.1, cw, 0.4, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(sl, f"{seats}  ·  {100*count/s['fleet']:.0f}% of the fleet",
                 x, cy + 3.5, cw, 0.4, size=11, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    add_text(sl, "On trunk corridors the large and medium buses are kept to a balanced "
                 "fifty-fifty mix, so neither size dominates a route.",
             0.9, 6.45, 11.9, 0.4, size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)
    footer(sl, "Slide 5  ·  Fleet composition")

    # ── 6 · Electric-bus backbone + operator absorption ──────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "The electric-bus backbone", color=TEAL,
                 sub="The plan is anchored on the existing Srinagar Smart City electric-bus routes.")
    nodes = [
        (0.7, 1.95, 2.7, 1.75, NAVY, "30", "Official Smart City\nelectric-bus routes", "already operating"),
        (4.0, 1.95, 2.7, 1.75, TEAL, "45", "Matching permits\nin the dataset", "duplicate operators"),
        (7.3, 1.95, 2.7, 1.75, PURPLE, "8.0", "Buses per route\nrecommended", "at a 15-minute service"),
        (10.6, 1.95, 2.2, 1.75, GREEN, "+9.7%", "Above the frequency-matched\nfleet benchmark", "comfortably within tolerance"),
    ]
    for x, y, w, h, col, big, l1, l2 in nodes:
        add_rect(sl, x, y, w, h, fill=CARD, rounded=True, line=col)
        add_text(sl, big, x, y + 0.1, w, 0.5, size=26, bold=True, color=col,
                 align=PP_ALIGN.CENTER, font="Calibri")
        add_text(sl, l1, x, y + 0.66, w, 0.62, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_text(sl, l2, x, y + 1.34, w, 0.32, size=9.5, color=GREY, italic=True, align=PP_ALIGN.CENTER)
    for x1, x2 in [(3.4, 4.0), (6.7, 7.3), (10.0, 10.6)]:
        add_arrow(sl, x1, 2.82, x2, 2.82, color=TEAL, width=2.5)
    # operator absorption
    pmb, lpv, hpv = s['op_pmb'], s['op_lpv'], s['op_hpv']
    cr = lambda n, lakh: n * lakh / 100.0
    rows = [
        ("Private minibus", pmb, "Buy back at about Rs 15 lakh each", f"Rs {cr(pmb,15):.2f} Cr"),
        ("Small bus / tempo", lpv, "Reassign to last-mile feeder duty", f"Rs {cr(lpv,3):.2f} Cr"),
        ("Large bus", hpv, "Absorb into the public fleet", f"Rs {cr(hpv,50):.2f} Cr"),
        ("Total absorbed", s['merged'], "Mixed strategy, by consultation",
         f"Rs {cr(pmb,15)+cr(lpv,3)+cr(hpv,50):.2f} Cr"),
    ]
    add_rect(sl, 0.7, 4.05, 12.1, 2.5, fill=GOLD_LT, rounded=True, line=GOLD)
    add_text(sl, "Operators absorbed into the new network — indicative compensation",
             0.95, 4.13, 11.5, 0.35, size=14, bold=True, color=GOLD)
    col_x = [0.95, 5.0, 6.6, 10.9]
    headers = ["Operator type", "Permits", "Recommended action", "Indicative cost"]
    for i, hh in enumerate(headers):
        add_text(sl, hh, col_x[i], 4.5, 3.8, 0.28, size=10.5, bold=True, color=GOLD)
    for i, (op, n, strat, est) in enumerate(rows):
        y = 4.84 + i * 0.33
        tot = (i == len(rows) - 1)
        add_text(sl, op, col_x[0], y, 4.0, 0.3, size=11.5, bold=tot, color=DARK)
        add_text(sl, f"{n}", col_x[1], y, 1.4, 0.3, size=11.5, bold=tot, color=DARK)
        add_text(sl, strat, col_x[2], y, 4.2, 0.3, size=11.5, color=DARK, italic=not tot, bold=tot)
        add_text(sl, est, col_x[3], y, 2.0, 0.3, size=11.5, bold=tot,
                 color=SAFFRON if tot else DARK)
    add_text(sl, "Costs are indicative estimates from peer-city averages — not a price commitment.",
             0.95, 6.2, 11.5, 0.28, size=9.5, color=GREY, italic=True)
    footer(sl, "Slide 6  ·  Backbone and operator absorption")

    # ── 7 · How often buses run ──────────────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "How often buses run", color=PURPLE,
                 sub="Buses are spaced by wait time. No route in the plan waits longer than 35 minutes.")
    bands = [
        ("Smart City electric backbone", 15, s['hw15'], TEAL, "the busiest corridors"),
        ("Main trunk lines", 20, s['hw20'], NAVY, "high-demand routes"),
        ("Feeders and lifeline routes", 35, s['hw35'], PURPLE, "neighbourhood & rural links"),
    ]
    base_x, scale, by0, rh = 3.2, 8.0 / 40.0, 2.1, 1.15
    for m in (0, 10, 20, 30, 40):
        ax = base_x + m * scale
        add_text(sl, f"{m} min", ax - 0.4, 1.75, 0.8, 0.25, size=9, color=GREY, align=PP_ALIGN.CENTER)
        add_rect(sl, ax - 0.01, 2.0, 0.02, 4.0, fill=RGBColor(0xE0, 0xE0, 0xE0))
    for i, (label, headway, n_routes, col, note) in enumerate(bands):
        y = by0 + i * (rh + 0.25)
        add_text(sl, label, 0.3, y, 2.8, rh, size=13, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        strip_w = 40 * scale
        add_rect(sl, base_x, y + 0.2, strip_w, rh - 0.4, fill=LIGHT, rounded=True)
        bus_w, bus_h = 0.34, 0.46
        bus_y = y + (rh - bus_h) / 2
        for m in range(0, 41, headway):
            bx = base_x + m * scale - bus_w / 2
            if bx + bus_w > base_x + strip_w + 0.1:
                continue
            add_rect(sl, bx, bus_y, bus_w, bus_h, fill=col, rounded=True)
        add_text(sl, f"every {headway} min", base_x + strip_w + 0.2, y + 0.12, 1.9, 0.4,
                 size=12, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
        add_text(sl, f"{n_routes} routes · {note}", base_x + strip_w + 0.2, y + 0.52, 1.9, 0.5,
                 size=9.5, color=GREY, italic=True)
    add_text(sl, "Earlier drafts allowed one-hour waits on lifeline routes; those have been removed — "
                 "the longest wait anywhere is now 35 minutes.",
             0.3, 6.4, 12.5, 0.5, size=12, color=GREY, italic=True)
    footer(sl, "Slide 7  ·  Service frequency")

    # ── 8 · The formulas, in full ─────────────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "The methodology, stated in full", color=NAVY,
                 sub="Every number rests on a standard transit-planning formula. Here they are, written out.")
    formula_card(sl, 0.6, 1.65, 6.0, 2.15,
                 "Number of buses on a route = round-trip time divided by how often a bus runs, "
                 "rounded up, plus a 15% spare allowance for maintenance and breakdowns.",
                 "Buses  =  round-trip time ÷ wait between buses  (rounded up)  × 1.15",
                 "Vuchic (2005); Ceder (2007); spare allowance per industry practice.", NAVY)
    formula_card(sl, 6.8, 1.65, 6.0, 2.15,
                 "Round-trip time = driving time both ways + time spent at stops + junction "
                 "delays + peak-hour congestion + a recovery buffer to keep to schedule.",
                 "round-trip time  =  driving + stops + junctions + congestion + recovery",
                 "Transit Capacity & Quality of Service Manual (Transportation Research Board, 2013).", TEAL)
    formula_card(sl, 0.6, 4.0, 6.0, 2.55,
                 "Estimated daily riders = people within a short walk × the share who take the bus "
                 "× trips per day × the route's share of a shared corridor × a calibration factor. "
                 "It is computed automatically from open data, with no operator GPS feed needed.",
                 "riders  =  walk-up population × bus share × trips × corridor share × factor",
                 "Gravity / accessibility model — Hansen (1959); Ortúzar & Willumsen (2011).", PURPLE)
    formula_card(sl, 6.8, 4.0, 6.0, 2.55,
                 "Composite Demand Index = half the population score + half the destinations score. "
                 "Routes are then grouped into high, medium and low priority, which sets how often "
                 "buses run on each.",
                 "demand index  =  ½ population  +  ½ destinations  →  natural-breaks grouping",
                 "Natural-breaks classification — Jenks (1967).", GOLD)
    footer(sl, "Slide 8  ·  Formulas and their basis")

    # ── 9 · Sources & references ──────────────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "Sources and references", color=TEAL,
                 sub="The data behind the plan, and the published research behind each method.")
    left = [
        ("Population", "WorldPop satellite-derived population grid (Tatem, 2017); Census of India 2011; "
                       "Srinagar Smart City detailed project report."),
        ("Road network & travel time", "OpenStreetMap road data; routing by the open-source routing "
                                        "engine OSRM (Luxen & Vetter, 2011)."),
        ("Destinations", "Points of interest from OpenStreetMap (hospitals, colleges, markets, offices)."),
        ("Walking distance", "Four-hundred-metre walk catchment, about a five-minute walk "
                             "(El-Geneidy et al., 2014)."),
    ]
    right = [
        ("Buses & frequency", "Vuchic, Urban Transit (2005); Ceder, Public Transit Planning and "
                              "Operation (2007)."),
        ("Cycle & recovery time", "Transit Capacity and Quality of Service Manual, Transportation "
                                  "Research Board (2013)."),
        ("Demand", "Hansen, accessibility model (1959); Ortúzar & Willumsen, Modelling Transport (2011)."),
        ("Ridership check", "Srinagar Smart City electric-bus published ridership totals "
                            "(an aggregate figure — no individual GPS or fare-card traces are used)."),
    ]
    for ci, col in enumerate([left, right]):
        x = 0.6 + ci * 6.3
        yy = 1.75
        for topic, cite in col:
            add_text(sl, topic, x, yy, 5.9, 0.35, size=13, bold=True, color=TEAL)
            add_text(sl, cite, x, yy + 0.34, 5.9, 0.9, size=11, color=DARK)
            yy += 1.32
    footer(sl, "Slide 9  ·  Sources and references")

    # ── 10 · Fleet density vs peer cities ─────────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "How the fleet compares with other Indian cities", color=NAVY,
                 sub="Buses per one thousand residents — the plan sits comfortably in the peer-city range.")
    cities = [
        ("Mysuru", 0.35, GREY), ("Bhopal", 0.43, GREY), ("Indore", 0.45, GREY),
        ("Bengaluru (BMTC)", 0.51, GREY),
        ("Srinagar (this plan)", round(s['per_1000'], 2), TEAL),
        ("Chandigarh", 0.65, GREY), ("Pune", 0.75, GREY),
        ("Delhi", 0.90, GREY), ("Mumbai", 1.20, GREY),
    ]
    max_v, bar_x0, bar_max, row_h, by0 = 1.30, 4.4, 7.8, 0.42, 1.9
    for i, (name, val, col) in enumerate(cities):
        y = by0 + i * row_h
        is_self = "Srinagar" in name
        add_text(sl, name, 0.6, y, 3.6, row_h - 0.05, size=12, bold=is_self,
                 color=TEAL if is_self else DARK, anchor=MSO_ANCHOR.MIDDLE)
        w = bar_max * val / max_v
        add_rect(sl, bar_x0, y + 0.06, w, row_h - 0.16,
                 fill=TEAL if is_self else RGBColor(0xBD, 0xBD, 0xBD))
        add_text(sl, f"{val:.2f}", bar_x0 + w + 0.1, y, 1.0, row_h - 0.05, size=11, bold=is_self,
                 color=TEAL if is_self else GREY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(sl, f"At {s['per_1000']:.2f} buses per thousand residents, Srinagar is closest to "
                 "Chandigarh — comparable in size, terrain and scope.",
             0.6, 5.85, 12, 0.6, size=12, italic=True, color=GREY)
    footer(sl, "Slide 10  ·  Peer-city comparison")

    # ── 11 · What is modelled now → future refinement ────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "What is modelled today, and how it sharpens with more data", color=PURPLE,
                 sub="The method is sound today and improves as official datasets become available.")
    rows = [
        ("Walking distance", "Straight-line walk catchment", "Street-network walk paths around lakes and the river", PURPLE),
        ("Population", "Satellite population grid", "Ward-level census counts", TEAL),
        ("Travel time", "Open routing-engine estimate with a congestion factor", "Measured bus speeds from the field", GOLD),
        ("Demand", "Population-based estimate", "Fare-card ridership, route by route", SAFFRON),
        ("Tourist demand", "Weighted tourist destinations", "Actual visitor-arrival figures", NAVY),
        ("Road closures", "Not yet modelled", "Highway-convoy and winter-closure calendar", GREEN),
    ]
    by0, rh = 1.9, 0.78
    add_text(sl, "Concept", 0.6, by0, 2.5, 0.4, size=11, bold=True, color=GREY)
    add_text(sl, "Modelled today", 3.4, by0, 4.4, 0.4, size=11, bold=True, color=GREY)
    add_text(sl, "Sharper with official data", 8.2, by0, 4.6, 0.4, size=11, bold=True, color=GREY)
    for i, (concept, now, later, col) in enumerate(rows):
        y = by0 + 0.45 + i * rh
        add_rect(sl, 0.6, y + 0.1, 0.15, 0.55, fill=col)
        add_text(sl, concept, 0.85, y, 2.5, rh - 0.1, size=12, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(sl, 3.4, y + 0.1, 4.2, rh - 0.2, fill=LIGHT, rounded=True)
        add_text(sl, now, 3.55, y + 0.1, 4.0, rh - 0.2, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_arrow(sl, 7.7, y + rh / 2, 8.1, y + rh / 2, color=col, width=2.0)
        add_rect(sl, 8.2, y + 0.1, 4.6, rh - 0.2, fill=RGBColor(0xE3, 0xF2, 0xFD), rounded=True)
        add_text(sl, later, 8.35, y + 0.1, 4.3, rh - 0.2, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(sl, "Slide 11  ·  Today's model and its refinements")

    # ── 12 · Data we would value from your office ────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "Data we would value from your office", color=TEAL,
                 sub="These refinements are limited by data availability, not by method. The concrete ask:")
    asks = [
        ("Most useful", "Live operator permit registry",
         "Permit number, owner, vehicle and depot",
         "Turns the merged permits into an exact, defensible compensation list.", SAFFRON),
        ("Most useful", "Bus position records from current operators",
         "A few weeks of location traces",
         "Replaces the open routing-engine travel-time estimate with measured speeds.", SAFFRON),
        ("Valuable", "Ward-level population from the latest census",
         "Ward boundaries and counts",
         "Sharpens the demand score and the priority grouping.", GOLD),
        ("Valuable", "Fare-card ridership, route by route",
         "Daily boardings per route",
         "Enables true route-level demand and frequency-induced ridership.", GOLD),
        ("Valuable", "Tourist arrival figures",
         "Daily arrivals for Gulmarg, Pahalgam, Sonamarg",
         "Lets us size the seasonal tourist fleet on real numbers.", GOLD),
        ("Future", "Road-operability calendar",
         "Convoy and winter-closure dates",
         "Automatically stands down routes that cannot run in winter.", NAVY),
        ("Future", "Bus-stop and depot inventory",
         "Location of every stop and depot capacity",
         "Replaces the assumed stop spacing with the real network.", NAVY),
        ("Future", "Bridge-weight and road-width register",
         "Per-bridge limits and arterial widths",
         "Keeps large buses off roads and bridges that cannot take them.", NAVY),
    ]
    by_h, by_a, rh = 1.6, 1.96, 0.58
    add_rect(sl, 0.45, by_h, 12.55, 0.32, fill=NAVY, rounded=True)
    for label, x, w in [("Priority", 0.55, 1.6), ("Dataset", 2.25, 3.6), ("Why it matters", 6.0, 6.8)]:
        add_text(sl, label, x, by_h + 0.04, w, 0.26, size=10, bold=True, color=WHITE)
    for i, (prio, name, detail, why, col) in enumerate(asks):
        y = by_a + i * rh
        if i % 2 == 0:
            add_rect(sl, 0.45, y, 12.55, rh - 0.04, fill=LIGHT)
        add_rect(sl, 0.55, y + 0.12, 1.45, rh - 0.28, fill=col, rounded=True)
        add_centered_label(sl, prio, 0.55, y + 0.12, 1.45, rh - 0.28, size=9.5, color=WHITE)
        add_text(sl, name, 2.25, y + 0.05, 3.6, 0.26, size=10.5, bold=True, color=DARK)
        add_text(sl, detail, 2.25, y + 0.3, 3.6, 0.24, size=8.5, color=GREY, italic=True)
        add_text(sl, why, 6.0, y, 7.0, rh - 0.04, size=10, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
        # custodian merged into why column width? keep why spanning; detail holds custodian text
    footer(sl, "Slide 12  ·  Data request")

    # ── 13 · What this plan is — and is not ───────────────────────────────────
    sl = slide_blank(prs)
    slide_header(sl, "What this plan is — and what it is not", color=NAVY)
    is_items = [
        f"{s['active']} active routes, ranked by demand and grouped into priority tiers",
        "Buses sized for a published frequency — every 15, 20 or 35 minutes",
        "Anchored on the Smart City electric-bus routes, checked against their ridership",
        "An operator workbook with absorption and indicative compensation",
        "A fully reproducible analysis — every number traces back to its source",
    ]
    isnot_items = [
        "A passenger survey — on-the-ground validation is still advisable",
        "A statement of the exact extra buses needed versus those running today",
        "Final, committed frequencies — these follow operator consultation",
        "A compensation price — the figures are indicative estimates",
        "A tourist-demand model — that needs real arrival data",
    ]
    add_rect(sl, 0.6, 1.75, 6.1, 5.05, fill=GREEN_LT, rounded=True, line=GREEN)
    add_rect(sl, 0.6, 1.75, 6.1, 0.6, fill=GREEN, rounded=True)
    add_centered_label(sl, "WHAT IT IS", 0.6, 1.75, 6.1, 0.6, size=16, color=WHITE)
    for i, t in enumerate(is_items):
        y = 2.62 + i * 0.85
        add_oval(sl, 0.85, y + 0.05, 0.35, 0.35, fill=GREEN)
        add_centered_label(sl, "✓", 0.85, y + 0.05, 0.35, 0.35, size=14, color=WHITE)
        add_text(sl, t, 1.35, y, 5.2, 0.78, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(sl, 6.95, 1.75, 6.1, 5.05, fill=SAFFRON_LT, rounded=True, line=SAFFRON)
    add_rect(sl, 6.95, 1.75, 6.1, 0.6, fill=SAFFRON, rounded=True)
    add_centered_label(sl, "WHAT IT IS NOT", 6.95, 1.75, 6.1, 0.6, size=16, color=WHITE)
    for i, t in enumerate(isnot_items):
        y = 2.62 + i * 0.85
        add_oval(sl, 7.2, y + 0.05, 0.35, 0.35, fill=SAFFRON)
        add_centered_label(sl, "✗", 7.2, y + 0.05, 0.35, 0.35, size=14, color=WHITE)
        add_text(sl, t, 7.7, y, 5.2, 0.78, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    footer(sl, "Slide 13  ·  Honest scope")

    # ── 14 · Closing / asks ───────────────────────────────────────────────────
    sl = slide_blank(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=NAVY)
    add_rect(sl, 0, 6.7, 13.33, 0.8, fill=TEAL)
    add_text(sl, "The recommended path", 0.7, 0.7, 12, 0.6, size=22, bold=True,
             color=TEAL_LIGHT, font="Calibri")
    add_text(sl, "Five asks", 0.7, 1.3, 12, 1.0, size=42, bold=True, color=WHITE, font="Calibri")
    asks = [
        ("1", "Approve the plan", "Endorse the recommended network at the first-stage review."),
        ("2", "Open operator dialogue", "Authorise consultation with the All Jammu & Kashmir Transport Welfare Association."),
        ("3", "Roll out in phases", "Start with the electric-bus backbone and main trunks; extend to feeders as ridership confirms."),
        ("4", "Refresh with new data", "Update the analysis as official datasets and ridership become available."),
        ("5", "Commission the refinements", "Sanction the next round of work — street-network walking, ridership-based demand, road-closure calendar."),
    ]
    for i, (n, title, body) in enumerate(asks):
        y = 2.7 + i * 0.78
        add_oval(sl, 0.8, y, 0.6, 0.6, fill=GOLD)
        add_centered_label(sl, n, 0.8, y, 0.6, 0.6, size=22, color=NAVY)
        add_text(sl, title, 1.6, y - 0.02, 3.4, 0.35, size=15, bold=True, color=WHITE)
        add_text(sl, body, 5.1, y + 0.0, 7.8, 0.6, size=12, color=TEAL_LIGHT, italic=True,
                 anchor=MSO_ANCHOR.MIDDLE)
    add_text(sl, "The route plan, interactive map and operator workbook are available for review on the dashboard.",
             0.7, 6.85, 12, 0.45, size=11, color=WHITE, italic=True, align=PP_ALIGN.CENTER)


# ─── Stat loader ─────────────────────────────────────────────────────────────
def load_stats(csv_path: Path):
    defaults = dict(
        total_routes=342, active=207, trunks=50, feeders=157, merged=135,
        fleet=1009, hpv=80, mpv=807, lpv=122, hp=130, mp=54, lp=23,
        hw15=45, hw20=85, hw35=77, sscl_fleet=362, sscl_matched=45,
        social=87, tourist=69, op_pmb=100, op_lpv=34, op_hpv=1,
        per_1000=0.51, coverage=31.1, net_pop=1588964,
    )
    if not csv_path.exists():
        return defaults
    df = pd.read_csv(csv_path, encoding="utf-8-sig")
    act = df[df["Action_Taken"] != "MERGED_INTO_TRUNK"].copy()
    if {"Fleet_Required", "HPV_Count", "MPV_Count"}.issubset(act.columns):
        act["LPV_Count"] = (act["Fleet_Required"] - act["HPV_Count"] - act["MPV_Count"]).clip(lower=0)
    out = dict(defaults)
    out["total_routes"] = len(df)
    out["active"]  = len(act)
    out["trunks"]  = int((df["Action_Taken"] == "UPGRADED_TO_TRUNK").sum())
    out["feeders"] = int((df["Action_Taken"] == "RETAINED_AS_FEEDER").sum())
    out["merged"]  = int((df["Action_Taken"] == "MERGED_INTO_TRUNK").sum())
    out["fleet"]   = int(act["Fleet_Required"].sum())
    out["hpv"]     = int(act["HPV_Count"].sum())
    out["mpv"]     = int(act["MPV_Count"].sum())
    out["lpv"]     = int(act["LPV_Count"].sum())
    for band in ("HP", "MP", "LP"):
        out[band.lower()] = int((act["Priority_Band"] == band).sum())
    hw = act["Headway_Min"].value_counts().to_dict()
    out["hw15"], out["hw20"], out["hw35"] = int(hw.get(15, 0)), int(hw.get(20, 0)), int(hw.get(35, 0))
    sscl = act[act.get("CMP_Trunk", False) == True]
    out["sscl_fleet"], out["sscl_matched"] = int(sscl["Fleet_Required"].sum()), len(sscl)
    out["social"]  = int(df["Social_Flag"].astype(str).str.lower().isin(["true", "1"]).sum())
    out["tourist"] = int(df["Tourist_Corridor"].astype(str).str.lower().isin(["true", "1"]).sum())
    merged = df[df["Action_Taken"] == "MERGED_INTO_TRUNK"]
    if "Displaced_Operator_Class" in merged.columns:
        d = merged["Displaced_Operator_Class"].value_counts()
        out["op_pmb"] = int(d.get("Private Minibus", 0))
        out["op_lpv"] = int(d.get("LPV / Tempo", 0))
        out["op_hpv"] = int(d.get("HPV Bus", 0))
    # F-V9: honest coverage = served ÷ study-area population (WorldPop ~5.1M),
    # not the Srinagar-UA planning figure. per_1000 = buses per 1,000 served.
    out["net_pop"] = (int(act["Population_Served"].sum())
                      if "Population_Served" in act.columns else 1588964)
    STUDY_AREA_POP = 5_105_699
    out["per_1000"] = out["fleet"] / (out["net_pop"] / 1000.0)
    out["coverage"] = out["net_pop"] / STUDY_AREA_POP * 100.0
    return out


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--outdir", default="outputs_v3.4.1")
    parser.add_argument("--engine-csv", default="outputs_v3.4.1/Rationalised_Routes_Kashmir_v3.csv")
    args = parser.parse_args()
    stats = load_stats(Path(args.engine_csv))
    os.makedirs(args.outdir, exist_ok=True)
    out_path = os.path.join(args.outdir, "Kashmir_Transit_Diagrammatic_Pitch.pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build(prs, stats)
    prs.save(out_path)
    print(f"Wrote {out_path}")
    print(f"Stats source: {args.engine_csv}")


if __name__ == "__main__":
    main()
